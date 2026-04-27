"""
Paper → Slides: 自动从论文生成演示文稿

核心功能：
1. 论文内容提取（标题、摘要、方法、实验、结论）
2. LLM 生成幻灯片结构
3. 多格式输出（PPTX/MD/HTML）
4. 演讲者备注生成
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional, List, Dict, Any

# 可选依赖：python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from db.database import Database
from sections.segment import segment_into_sections, format_section_snippets


@dataclass
class SlidesConfig:
    """幻灯片生成配置."""
    template: str = "academic"       # academic | minimal | modern
    num_slides: int = 10            # 幻灯片数量
    output_format: str = "pptx"      # pptx | md | html
    output_path: Optional[str] = None
    include_notes: bool = False      # 包含演讲者备注
    language: str = "zh"             # zh | en | bilingual


@dataclass
class Slide:
    """单个幻灯片."""
    title: str
    content: str
    notes: str = ""
    slide_type: str = "content"  # title | content | comparison | summary


@dataclass
class SlidesResult:
    """生成结果."""
    output_path: str
    slide_count: int
    paper_count: int
    slides: List[Slide] = field(default_factory=list)


class PaperSlidesGenerator:
    """论文到幻灯片生成器."""

    # 幻灯片模板结构
    TEMPLATES = {
        "academic": {
            "title_slide": {"layout": "title", "bg_color": None},
            "content_slide": {"layout": "content", "font_title": 32, "font_body": 18},
            "section_slide": {"layout": "section_header"},
        },
        "minimal": {
            "title_slide": {"layout": "blank"},
            "content_slide": {"layout": "blank", "font_title": 28, "font_body": 16},
            "section_slide": {"layout": "blank"},
        },
        "modern": {
            "title_slide": {"layout": "title", "bg_color": (0, 100, 180)},
            "content_slide": {"layout": "content", "font_title": 36, "font_body": 20},
            "section_slide": {"layout": "section_header", "bg_color": (240, 240, 240)},
        },
    }

    def __init__(self, db: Optional[Database] = None):
        self.db = db or Database()
        self._llm_client = None

    @property
    def llm_client(self):
        """延迟加载 LLM 客户端."""
        if self._llm_client is None:
            from llm.client import get_default_client
            self._llm_client = get_default_client()
        return self._llm_client

    def generate(
        self,
        paper_ids: List[str],
        config: Optional[SlidesConfig] = None,
    ) -> SlidesResult:
        """生成幻灯片.

        Args:
            paper_ids: 论文 ID 列表
            config: 生成配置

        Returns:
            SlidesResult: 生成结果
        """
        config = config or SlidesConfig()

        # 1. 获取论文内容
        papers_content = self._fetch_papers_content(paper_ids)

        # 2. 生成幻灯片结构
        slides = self._generate_slide_structure(papers_content, config)

        # 3. 输出
        output_path = self._write_output(slides, config)

        return SlidesResult(
            output_path=str(output_path),
            slide_count=len(slides),
            paper_count=len(paper_ids),
            slides=slides,
        )

    def _fetch_papers_content(self, paper_ids: List[str]) -> List[Dict[str, Any]]:
        """获取论文内容."""
        papers = []
        for pid in paper_ids:
            paper = self.db.get_paper(pid)
            if not paper:
                continue

            # 提取关键内容
            content = {
                "id": pid,
                "title": paper.get("title", ""),
                "authors": paper.get("authors", ""),
                "abstract": paper.get("abstract", ""),
                "year": paper.get("published", "")[:4] if paper.get("published") else "",
                "tags": paper.get("tags", []),
                "plain_text": paper.get("plain_text", ""),
            }

            # 如果有全文，提取关键章节
            if content["plain_text"]:
                sections = segment_into_sections(content["plain_text"])
                content["sections"] = sections
                content["snippet"] = format_section_snippets(sections)
            else:
                content["sections"] = []
                content["snippet"] = content["abstract"]

            papers.append(content)

        return papers

    def _generate_slide_structure(
        self,
        papers: List[Dict[str, Any]],
        config: SlidesConfig,
    ) -> List[Slide]:
        """生成幻灯片结构.

        核心算法：
        1. 构建幻灯片大纲（标题页 → 背景 → 方法 → 实验 → 结论）
        2. 从论文内容提取引用
        3. 生成演讲者备注
        """
        if not papers:
            return [Slide(title="No Content", content="No papers found")]

        # 单论文 vs 多论文
        if len(papers) == 1:
            return self._generate_single_paper_slides(papers[0], config)
        else:
            return self._generate_comparison_slides(papers, config)

    def _generate_single_paper_slides(
        self,
        paper: Dict[str, Any],
        config: SlidesConfig,
    ) -> List[Slide]:
        """单论文幻灯片生成."""
        slides = []

        # 1. 标题页
        slides.append(Slide(
            title=paper["title"],
            content=f"{paper['authors']}\n{paper['year']}",
            notes="开场介绍论文标题和作者",
            slide_type="title",
        ))

        # 2. 摘要/动机
        abstract = paper.get("abstract", "")[:500]
        slides.append(Slide(
            title="研究动机",
            content=abstract,
            notes="介绍研究背景和动机，强调问题的重要性",
            slide_type="content",
        ))

        # 3. 关键方法
        sections = paper.get("sections", [])
        method_sections = [s for s in sections if any(
            kw in s[0].lower() for kw in ["method", "approach", "model", "architecture"]
        )]

        if method_sections:
            for title, content in method_sections[:2]:
                slides.append(Slide(
                    title=f"方法: {title}",
                    content=content[:800],
                    notes=f"详细讲解{title}部分的技术细节",
                    slide_type="content",
                ))

        # 4. 实验结果
        results_sections = [s for s in sections if any(
            kw in s[0].lower() for kw in ["experiment", "result", "evaluation"]
        )]

        if results_sections:
            for title, content in results_sections[:1]:
                slides.append(Slide(
                    title="实验结果",
                    content=content[:600],
                    notes="展示关键实验数据和方法对比",
                    slide_type="content",
                ))

        # 5. 结论
        conclusion_sections = [s for s in sections if "conclusion" in s[0].lower()]
        if conclusion_sections:
            title, content = conclusion_sections[0]
            slides.append(Slide(
                title="结论",
                content=content[:500],
                notes="总结论文贡献和未来工作方向",
                slide_type="summary",
            ))

        # 6. 限制/未来工作
        slides.append(Slide(
            title="参考与引用",
            content=f"Tags: {', '.join(paper.get('tags', []))}",
            notes="提供进一步阅读的建议",
            slide_type="content",
        ))

        return slides[:config.num_slides]

    def _generate_comparison_slides(
        self,
        papers: List[Dict[str, Any]],
        config: SlidesConfig,
    ) -> List[Slide]:
        """多论文对比幻灯片."""
        slides = []

        # 标题页
        titles = [p["title"][:40] for p in papers]
        slides.append(Slide(
            title="论文对比分析",
            content="\n".join(f"• {t}" for t in titles),
            notes="介绍即将对比的论文",
            slide_type="title",
        ))

        # 对比表格
        comparison = self._generate_comparison_table(papers)
        slides.append(Slide(
            title="论文概览对比",
            content=comparison,
            notes="展示各论文基本信息",
            slide_type="comparison",
        ))

        # 逐个论文简介
        for paper in papers:
            slides.append(Slide(
                title=paper["title"][:50],
                content=f"年份: {paper['year']}\n\n{paper.get('abstract', '')[:400]}",
                notes=f"介绍{paper['title']}的核心内容",
                slide_type="content",
            ))

        return slides[:config.num_slides]

    def _generate_comparison_table(self, papers: List[Dict[str, Any]]) -> str:
        """生成对比表格（Markdown 格式）."""
        headers = ["论文", "年份", "标签"]
        rows = []
        for p in papers:
            rows.append([
                p["title"][:30],
                p["year"],
                ", ".join(p.get("tags", [])[:3]),
            ])

        col_widths = [max(len(str(row[i])) for row in [headers] + rows) + 2
                      for i in range(len(headers))]

        lines = []
        # 表头
        header_line = " | ".join(h.ljust(w) for h, w in zip(headers, col_widths))
        lines.append(header_line)
        lines.append("|" + "|".join("-" * w for w in col_widths) + "|")

        # 数据行
        for row in rows:
            lines.append(" | ".join(str(cell).ljust(w) for cell, w in zip(row, col_widths)))

        return "\n".join(lines)

    def _write_output(
        self,
        slides: List[Slide],
        config: SlidesConfig,
    ) -> Path:
        """输出幻灯片文件."""
        output_path = config.output_path
        if not output_path:
            output_dir = Path.cwd() / "slides_output"
            output_dir.mkdir(exist_ok=True)
            timestamp = Path().name.replace(":", "-")
            ext = config.output_format
            output_path = str(output_dir / f"slides_{timestamp}.{ext}")

        if config.output_format == "md":
            return self._write_markdown(slides, output_path, config)
        elif config.output_format == "html":
            return self._write_html(slides, output_path, config)
        elif config.output_format == "pptx":
            return self._write_pptx(slides, output_path, config)
        else:
            return self._write_markdown(slides, output_path, config)

    def _write_markdown(
        self,
        slides: List[Slide],
        output_path: str,
        config: SlidesConfig,
    ) -> Path:
        """输出 Markdown 格式."""
        lines = []
        for i, slide in enumerate(slides, 1):
            lines.append(f"# Slide {i}: {slide.title}")
            lines.append("")
            if slide.slide_type == "title":
                lines.append(f"## {slide.content}")
            else:
                lines.append(slide.content)
            lines.append("")

            if config.include_notes and slide.notes:
                lines.append(f"**演讲备注**: {slide.notes}")
                lines.append("")

            lines.append("---")
            lines.append("")

        path = Path(output_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text("\n".join(lines), encoding="utf-8")
        return path

    def _write_html(
        self,
        slides: List[Slide],
        output_path: str,
        config: SlidesConfig,
    ) -> Path:
        """输出 HTML 格式."""
        slide_htmls = []
        for i, slide in enumerate(slides, 1):
            notes_html = f'<div class="notes">{slide.notes}</div>' if config.include_notes else ""
            slide_html = f"""
<div class="slide" id="slide-{i}">
    <h1>{slide.title}</h1>
    <div class="content">
        <pre>{slide.content}</pre>
    </div>
    {notes_html}
</div>
"""
            slide_htmls.append(slide_html)

        html = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>Paper Slides</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 0; padding: 20px; }}
        .slide {{ page-break-after: always; margin-bottom: 40px; border: 1px solid #ccc; padding: 20px; }}
        h1 {{ color: #333; border-bottom: 2px solid #0066cc; }}
        .content pre {{ white-space: pre-wrap; font-family: inherit; }}
        .notes {{ background: #f0f0f0; padding: 10px; margin-top: 20px; font-style: italic; }}
    </style>
</head>
<body>
    {''.join(slide_htmls)}
</body>
</html>
"""
        path = Path(output_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(html, encoding="utf-8")
        return path

    def _write_pptx(
        self,
        slides: List[Slide],
        output_path: str,
        config: SlidesConfig,
    ) -> Path:
        """输出 PPTX 格式（需要 python-pptx）."""
        if not HAS_PPTX:
            print("Warning: python-pptx 未安装，将输出 Markdown 格式")
            return self._write_markdown(slides, output_path.replace(".pptx", ".md"), config)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for slide_data in slides:
            if slide_data.slide_type == "title":
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
            else:
                slide = prs.slides.add_slide(prs.slide_layouts[6])

            # 添加标题
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5), Inches(12), Inches(1)
            )
            title_frame = title_box.text_frame
            title_frame.text = slide_data.title
            for paragraph in title_frame.paragraphs:
                paragraph.font.size = Pt(32)
                paragraph.font.bold = True

            # 添加内容
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.8), Inches(12), Inches(5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.text = slide_data.content[:2000]  # 限制长度

            # 添加备注
            if config.include_notes and slide_data.notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_data.notes

        path = Path(output_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(path))
        return path


# CLI 辅助函数
def main(argv=None):
    """CLI 入口点."""
    import sys
    sys.path.insert(0, str(Path(__file__).parent.parent))

    from cli.cmd.slides import slides
    slides.main(standalone_mode=False)


if __name__ == "__main__":
    main()
